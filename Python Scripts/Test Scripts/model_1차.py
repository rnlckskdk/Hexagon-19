import openai
from PyPDF2 import PdfReader
import docx2txt
import os 
import olefile
from pptx import Presentation

# 키 넣기

folder_path = r'C:\Users\82107\Desktop\코딩_공부\ChatGPT'
completions = []

# 모든 하위폴더를 탐색하는 코드 (dirpath, dirnames, filenames)
def explore_the_path(folder_path):
    ai_result = ""
    temp_paths = []
    for (root, dirs, files) in os.walk(folder_path):
        for file in files:
            temp_paths.append(os.path.join(root, file))
            if os.path.getsize(os.path.join(root,file)) < 4 * 1024 * 1024:  # 파일 크기가 4MB 이하인지 확인
                if file.endswith('txt'):
                    f = open(os.path.join(root,file), encoding='utf-8')
                    RD = f.read()                   
                    ai_result = openaiAPI(RD)
                    keywords_text = ai_result.choices[0].message.content
                    keywords_list = [line.split(". ", 1)[1].strip() for line in keywords_text.strip().split("\n") if ". " in line]
                    print("폴더경로 : ",root, "\n","파일이름 : ",file, "\n", keywords_list, sep="", end="\n\n")
                    
                elif file.endswith('pdf'):
                    RD = PdfReader(os.path.join(root,file))                     
                    new_pdf_text = RD.pages 
                    page_data = ""
                    for i in new_pdf_text:
                        one_page = i.extract_text()
                        page_data += one_page                    
                    ai_result = openaiAPI(page_data) 
                    keywords_text = ai_result.choices[0].message.content
                    keywords_list = [line.split(". ", 1)[1].strip() for line in keywords_text.strip().split("\n") if ". " in line]
                    print("폴더경로 : ",root, "\n","파일이름 : ",file, "\n", keywords_list, sep="", end="\n\n")
                    
                elif file.endswith('docx') and not file.startswith('~$'):
                    text = docx2txt.process(os.path.join(root,file))                    
                    ai_result = openaiAPI(text)
                    keywords_text = ai_result.choices[0].message.content
                    keywords_list = [line.split(". ", 1)[1].strip() for line in keywords_text.strip().split("\n") if ". " in line]
                    print("폴더경로 : ",root, "\n","파일이름 : ",file, "\n", keywords_list, sep="", end="\n\n")
                
                elif file.endswith('hwp'):
                    RD = olefile.OleFileIO(os.path.join(root,file))
                    encoded_text = RD.openstream('PrvText').read()
                    decoded_text = encoded_text.decode('utf-16')
                    ai_result = openaiAPI(decoded_text)
                    keywords_text = ai_result.choices[0].message.content
                    keywords_list = [line.split(". ", 1)[1].strip() for line in keywords_text.strip().split("\n") if ". " in line]
                    print("폴더경로 : ",root, "\n","파일이름 : ",file, "\n", keywords_list, sep="", end="\n\n")

                elif file.endswith('pptx'):
                    RD = Presentation(os.path.join(root,file))
                    text_runs = []
                    for slide in RD.slides:
                        for shape in slide.shapes:
                            if not shape.has_text_frame:
                                continue
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    text_runs.append(run.text)
                    combined_text = "\n".join(text_runs)
                    ai_result = openaiAPI(combined_text)
                    keywords_text = ai_result.choices[0].message.content
                    keywords_list = [line.split(". ", 1)[1].strip() for line in keywords_text.strip().split("\n") if ". " in line]
                    print("폴더경로 : ",root, "\n","파일이름 : ",file, "\n", keywords_list, sep="", end="\n\n")


def openaiAPI(data): 
    response = openai.chat.completions.create(
        model = 'gpt-3.5-turbo',
        messages = [
            {"role": "system", "content": "You are a sophisticated keyword extraction system."},
            {"role": "user", "content": "Please Extract 10 relevant keywords from the following text, numbered one by one."},
            {"role": "user", "content": data }
        ],
        temperature=0, # 낮은 값 :예측 가능성, 무작위성, 낮을수록 일관된 텍스트 
        top_p=0,  # 높은 top_p 값 : 창의적이고 다양한 응답, 덜 일관성, 낮은 top_p 값 : 일관된 결과, 창의성이 떨어질 수 있음
        )
    return response

explore_the_path(folder_path)