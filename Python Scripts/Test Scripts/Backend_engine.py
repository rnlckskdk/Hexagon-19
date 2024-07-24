import os
from PyPDF2 import PdfReader # pip install PyPDF2
import docx2txt # pip install docx2txt
import openai
# import time
import olefile
from pptx import Presentation
# import Database


#def get_Word_list(folder_path):    
#    return [f for f in os.listdir(folder_path) if f.endswith('.docx') and not f.startswith('~$')]

# open ai api 이용

# 여기에 키 들어감
openai.api_key = ""

def openaiAPI(data): 
    response = openai.chat.completions.create(
        model = 'gpt-3.5-turbo',
        messages = [
            {"role": "system", "content": "You are a sophisticated keyword extraction system."},
            {"role": "user", "content": "Please Extract 10 relevant keywords from the following text, numbered one by one."},
            {"role": "user", "content": data }
        ],
        temperature=0, # 예측 가능성, 무작위성 , 낮을수록 결정록적이고 일관된 텍스트 
        top_p=0,      # 다양성 , 상위
    )
    return response


def insertDB(readData, db, root, file, size):
    ai_result = openaiAPI(readData)
    keywords_text = ai_result.choices[0].message.content
    keywords_list = [line.split(". ", 1)[1].strip() for line in keywords_text.strip().split("\n") if ". " in line]
    db.insertFileInfo(root, file, keywords_list, str(size))
    # print(root, file, keywords_list)


# 모든 하위폴더를 탐색하는 코드
# 사용 결과 => dirpath, dirnames, filenames
# 파일의 위치를 정확히 알아야 하는가? yes
# 무엇이 필요한가? 
# 확장자명, 
# completions 리스트 사용
def explore_the_path(folder_path, db):
    print("in path")
    ai_result = ""
    temp_paths = []
    for (root, dirs, files) in os.walk(folder_path):
        for file in files:
            temp_paths.append(os.path.join(root, file))
            size = os.path.getsize(os.path.join(root,file))
            if size < 4 * 1024 * 1024:  # 파일 크기가 4MB 이하인지 확인
                if file.endswith('txt'):
                    f = open(os.path.join(root,file), encoding='utf-8')
                    RD = f.read() # read
                    # txt에서 text 추출 완료
                    insertDB(RD, db, root, file, size)
                    # print("파일경로 : ",root, "\n","파일이름 : ",file, "\n", ai_result.choices[0].message.content, sep="", end="\n\n")

                elif file.endswith('pdf'):
                    RD = PdfReader(os.path.join(root,file)) # PdfReader로 파일 읽어오기
                    new_pdf_text = RD.pages 
                    page_data = ""
                    for i in new_pdf_text:
                        one_page = i.extract_text()
                        page_data += one_page
                    # pdf에서 text 추출 완료
                    insertDB(page_data, db, root, file, size)
                    # print("파일경로 : ",root, "\n","파일이름 : ",file, "\n", ai_result.choices[0].message.content, sep="", end="\n\n")

                elif file.endswith('docx') and not file.startswith('~$'):
                    text = docx2txt.process(os.path.join(root,file))
                    # word에서 text 추출 완료
                    insertDB(text, db, root, file, size)
                    # print("파일경로 : ",root, "\n","파일이름 : ",file, "\n", ai_result.choices[0].message.content, sep="", end="\n\n")

                elif file.endswith('hwp'):
                    RD = olefile.OleFileIO(os.path.join(root,file))
                    encoded_text = RD.openstream('PrvText').read()
                    decoded_text = encoded_text.decode('utf-16')
                    insertDB(decoded_text, db, root, file, size)

                elif file.endswith('pptx'):
                    RD = Presentation(os.path.join(root,file))
                    text_runs = []
                    for slide in RD.slides:
                        for shape in slide.shapes:
                            if not shape.has_text_frame:
                                continue
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    text_runs.append(run.text)
                    combined_text = "\n".join(text_runs)
                    insertDB(combined_text, db, root, file, size)
                    # print("폴더경로 : ",root, "\n","파일이름 : ",file, "\n", ai_result.choices[0].message.content, sep="", end="\n\n")
    
    print("out path")
    