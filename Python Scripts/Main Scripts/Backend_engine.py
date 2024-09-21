import os
from PyPDF2 import PdfReader # pip install PyPDF2
import docx2txt # pip install docx2txt
import openai
# import time
import olefile
import tiktoken # 토큰 인코드/디코드 # pip install tiktoken 
from pptx import Presentation
# import Database

# open ai api 이용

# 여기에 키 들어감
openai.api_key = ""

encoder = tiktoken.encoding_for_model("gpt-3.5-turbo") # 토큰 활용
MAX_TOKENS = 8000 # API 사용 시 한 번에 보내고 받을 수 있는 최대 토큰 수

def openaiAPI(data): 
    response = openai.chat.completions.create(
        model = 'gpt-3.5-turbo',
        messages = [
            {"role": "system", "content": "You are a sophisticated keyword extraction system."},
            {"role": "user", "content": "Please Extract relevant keywords from the following text, regardless of the capacity, extract only 10 tags numbered one by one."},
            {"role": "user", "content": data }
        ],
        temperature=0, # 예측 가능성, 무작위성 , 낮을수록 결정록적이고 일관된 텍스트 
        top_p=0,      # 다양성 , 상위
    )
    return response


def insertDB(readData, db, root, file, size):
    tokens = encoder.encode(readData)

    # 큰 파일 처리
    if len(tokens) > MAX_TOKENS:
        split_readData = list(splitToken(tokens))
        keywords_list = []

        for readData_chunk in split_readData:
            decoded_chunk = encoder.decode(readData_chunk)
            ai_result = openaiAPI(decoded_chunk)  # 분할된 텍스트에 대해 API 호출
            keywords_text = ai_result.choices[0].message.content
            keywords_list.extend([line.split(". ", 1)[1].strip() for line in keywords_text.strip().split("\n") if ". " in line])

        db.insertFileInfo(root, file, keywords_list, str(size))
    else: # 기본
        ai_result = openaiAPI(readData)
        keywords_text = ai_result.choices[0].message.content
        keywords_list = [line.split(". ", 1)[1].strip() for line in keywords_text.strip().split("\n") if ". " in line]
        db.insertFileInfo(root, file, keywords_list, str(size))
        # print(root, file, keywords_list)

# 토큰 쪼개기
def splitToken(tokens, max_tokens=MAX_TOKENS):
    for i in range(0, len(tokens), max_tokens):
        yield tokens[i:i + max_tokens - 1]

# 인코딩 확인(아직 추가 안함)
def checkEncoding(readData, encoding='utf-8'):
    try:
        # 텍스트를 주어진 인코딩으로 변환
        readData.encode(encoding, errors='strict')
    except UnicodeEncodeError as e:
        # 인코딩할 수 없는 문자가 있는 경우 예외 처리
        print(f"error: {e}")

def explore_the_path(folder_path, db):
    print("in path")
    ai_result = ""
    temp_paths = []
    for (root, dirs, files) in os.walk(folder_path):
        for file in files:
            temp_paths.append(os.path.join(root, file))
            size = os.path.getsize(os.path.join(root,file))
            if file.endswith('txt'):
                f = open(os.path.join(root,file), encoding='utf-8')
                RD = f.read() # read
                # txt에서 text 추출 완료
                insertDB(RD, db, root, file, size)
                # print("파일경로 : ",root, "\n","파일이름 : ",file, "\n", ai_result.choices[0].message.content, sep="", end="\n\n")

            elif file.endswith('pdf'):
                RD = PdfReader(os.path.join(root,file)) # PdfReader로 파일 읽어오기
                new_pdf_text = RD.pages 
                page_data = ""
                for i in new_pdf_text:
                    one_page = i.extract_text()
                    page_data += one_page
                # pdf에서 text 추출 완료
                insertDB(page_data, db, root, file, size)
                # print("파일경로 : ",root, "\n","파일이름 : ",file, "\n", ai_result.choices[0].message.content, sep="", end="\n\n")

            elif file.endswith('docx') and not file.startswith('~$'):
                text = docx2txt.process(os.path.join(root,file))
                # word에서 text 추출 완료
                insertDB(text, db, root, file, size)
                # print("파일경로 : ",root, "\n","파일이름 : ",file, "\n", ai_result.choices[0].message.content, sep="", end="\n\n")

            elif file.endswith('hwp'):
                RD = olefile.OleFileIO(os.path.join(root,file))
                encoded_text = RD.openstream('PrvText').read()
                decoded_text = encoded_text.decode('utf-16')
                insertDB(decoded_text, db, root, file, size)

            elif file.endswith('pptx'):
                RD = Presentation(os.path.join(root,file))
                text_runs = []
                for slide in RD.slides:
                    for shape in slide.shapes:
                        if not shape.has_text_frame:
                            continue
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                text_runs.append(run.text)
                combined_text = "\n".join(text_runs)
                insertDB(combined_text, db, root, file, size)
                # print("폴더경로 : ",root, "\n","파일이름 : ",file, "\n", ai_result.choices[0].message.content, sep="", end="\n\n")
    print("out path")
    